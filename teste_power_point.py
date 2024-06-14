import os
from pptx import Presentation
from pptx.util import Inches, Pt
import platform
from collections import defaultdict

# Detecta o sistema operacional
sistema_operacional = platform.system()

# Define o separador de diretório com base no sistema operacional
versionador = os.sep


# Cria uma apresentação vazia
current_directory = os.getcwd()

prs = Presentation()

# Adiciona um slide de título
slide_layout = prs.slide_layouts[0]  # Layout 0 é um slide de título
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Título da Apresentação"
subtitle.text = "Subtítulo ou Descrição da Apresentação"

# Aumenta o tamanho da fonte do título
title_shape = title.text_frame
for paragraph in title_shape.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(32)  # Tamanho da fonte do título

# Aumenta o tamanho da fonte do subtítulo
subtitle_shape = subtitle.text_frame
for paragraph in subtitle_shape.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(24)  # Tamanho da fonte do subtítulo

# Define o tamanho das imagens em pixels (exemplo: 300x200 pixels)
img_width_px = 320
img_height_px = 320

# Converte o tamanho das imagens de pixels para polegadas
img_width = Pt(img_width_px / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
img_height = Pt(img_height_px / 96 * 72)

# Define a posição inicial das imagens e o espaço entre elas
margin_left = Inches(0.05)
margin_top = Inches(0.8)
horizontal_spacing = Inches(0.05)
vertical_spacing = Inches(0.05)

# Atualize este caminho conforme necessário
image_folder = 'C:\\Users\\eduardo.neto\\Desktop\\programa_v6\\graficos_gerados\\Graficos Single Pulse'

# Lista todas as imagens na pasta especificada
images = [f for f in os.listdir(image_folder) if f.endswith(('png', 'jpg', 'jpeg'))]

# Agrupa imagens por prefixo (considerando os primeiros n caracteres como prefixo)
n = 8  # Ajuste conforme necessário
grupos = defaultdict(list)
for image in images:
    prefix = image[:n]
    grupos[prefix].append(image)

# Adiciona slides para cada grupo de gráficos
for prefix, graficos in grupos.items():
    # Adiciona um slide vazio
    slide_layout = prs.slide_layouts[10]  # Layout 6 é um slide vazio
    slide = prs.slides.add_slide(slide_layout)

    # Adiciona título ao slide
    title_shape = slide.shapes.title
    if title_shape is None:
        # Se não houver espaço reservado para título, crie um textbox para o título
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        text_frame = title_shape.text_frame
        text_frame.text = f"Grupo: {prefix}"
    else:
        title_shape.text = f"Grupo: {prefix}"
        text_frame = title_shape.text_frame

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)  # Tamanho da fonte do título do slide

    for idx, grafico in enumerate(graficos):
        # Calcula a posição da imagem
        col = idx % 3  # Colunas (0, 1, 2)
        row = idx // 3  # Linhas (0, 1)
        left = margin_left + col * (img_width + horizontal_spacing)
        top = margin_top + row * (img_height + vertical_spacing)
        
        # Adiciona a imagem ao slide
        img_path = os.path.join(image_folder, grafico)
        print(image_folder)
        print(grafico)
        if os.path.exists(img_path):  # Verifica se o arquivo existe
            slide.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)
        else:
            print(f"Arquivo não encontrado: {img_path}")



prs.save(current_directory+versionador+'Template'+versionador+'output.pptx')
