import os
from pptx import Presentation
from pptx.util import Inches, Pt
import platform
from collections import defaultdict
import pandas as pd

sistema_operacional = platform.system()


if sistema_operacional == "Windows":
    versionador = '\\'
elif sistema_operacional == "Linux":
    versionador = '/'
elif sistema_operacional == "Darwin":
    print("Você está usando o macOS.")
else:
    print(f"Você está usando um sistema operacional desconhecido: {sistema_operacional}")
  
    


# Detecta o sistema operacional
sistema_operacional = platform.system()

# Define o separador de diretório com base no sistema operacional
versionador = os.sep

# Função para criar a apresentação de slides
def power_point(output_pptx):
    # Cria uma apresentação vazia
    current_directory = os.getcwd()

    prs = Presentation()
    
    
    
    
    
    
    
    
    # Adiciona um slide de título
    slide_layout = prs.slide_layouts[0]  # Layout 0 é um slide de título
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Single Pulse"
    subtitle.text = "Dados e Gráficos"
    
    # Aumenta o tamanho da fonte do título
    title_shape = title.text_frame
    for paragraph in title_shape.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(90)  # Tamanho da fonte do título
    
    # Aumenta o tamanho da fonte do subtítulo
    subtitle_shape = subtitle.text_frame
    for paragraph in subtitle_shape.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(52)  # Tamanho da fonte do subtítulo
    
    # Define o tamanho das imagens em pixels (exemplo: 300x200 pixels)
    
    
    img_width_px_1 = 320
    img_height_px_1 = 320
    
    img_width_px_2 = 320
    img_height_px_2 = 320
    
    # Converte o tamanho das imagens de pixels para polegadas
    img_width_1 = Pt(img_width_px_1 / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_1 = Pt(img_height_px_1 / 96 * 72)
    
    img_width_2 = Pt(img_width_px_2 / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_2 = Pt(img_height_px_2 / 96 * 72)
    
    # Define a posição inicial das imagens e o espaço entre elas
    margin_left_1 = Inches(1.4)
    margin_top_1 = Inches(0.8)
    horizontal_spacing_1 = Inches(0.01)
    vertical_spacing_1 = Inches(0.000)
    
    
    margin_left_2 = Inches(0.1)
    margin_top_2 = Inches(0.8)
    horizontal_spacing_2 = Inches(0.01)
    vertical_spacing_2 = Inches(0.01)
    
    
    
    # Define o caminho para a imagem
    img_caminho = 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\cnpem.png'
    
    
    img_width_px_img = 162
    img_height_px_img = 85
    
    # Converte o tamanho das imagens de pixels para polegadas
    img_width_img = Pt(img_width_px_img / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_img = Pt(img_height_px_img / 96 * 72)
    
    # Define as dimensões e a posição da imagem
    left_img = Inches(8)    # Distância da borda esquerda do slide
    top_img = Inches(0.01)     # Distância da borda superior do slide
    
    
    # Adiciona a imagem ao slide
    pic = slide.shapes.add_picture(img_caminho, left_img, top_img, img_width_img, img_height_img)
    
    # Opcional: Ajustar o tamanho da imagem sem alterar as proporções
    # pic = slide.shapes.add_picture(img_path, left, top)
    # pic.width = Inches(4)  # Define a largura
    # pic.height = Inches(3)  # Define a altura
    
    
    
    
    
    
    
    # Atualize este caminho conforme necessário
    image_folder = os.path.join(current_directory, 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\graficos_gerados\\Graficos Single Pulse')
    
    # Lista todas as imagens na pasta especificada
    images = [f for f in os.listdir(image_folder) if f.endswith(('png', 'jpg', 'jpeg'))]
    
    
    
    
    
    # Lista de arquivos na pasta
    files = [f for f in os.listdir(image_folder) if os.path.isfile(os.path.join(image_folder, f))]
    
    # Criação do DataFrame
    df = pd.DataFrame(files, columns=['Filename'])
    
    
    
    # Extraindo o prefixo (parte antes do '_') para usar como chave de agrupamento
    df['Group'] = df['Filename'].apply(lambda x: x.split('_')[0])
    
    
    # Agrupando por 'Group'
    grouped = df.groupby('Group')
    # Realizando o agrupamento
    
    
    # Criando um único dicionário com todos os grupos
    combined_dict = {key: group.to_dict('list') for key, group in grouped}
    
    
        
    
    
    
    
    # Criando um novo dicionário plano
    novo_dicionario = {}
    
    for chave_principal, valor in combined_dict.items():
        novo_dicionario[chave_principal] = valor['Filename']
    
    # Imprimindo o novo dicionário
    del novo_dicionario['Grafico Single Pulse.opju']
    print(novo_dicionario)     
    
    
    
    
    # Adiciona slides para cada grupo de gráficos
    for prefix, graficos in novo_dicionario.items():
        # Adiciona um slide vazio
        slide_layout = prs.slide_layouts[6]  # Layout 6 é um slide vazio
        slide1 = prs.slides.add_slide(slide_layout)
        slide2 = prs.slides.add_slide(slide_layout)
        pic = slide2.shapes.add_picture(img_caminho, left_img, top_img, img_width_img, img_height_img)
        pic = slide1.shapes.add_picture(img_caminho, left_img, top_img, img_width_img, img_height_img)
    
        # Adiciona título ao slide
        title_shape1 = slide1.shapes.title
        title_shape2 = slide2.shapes.title
    
        if (title_shape1 is None) or (title_shape2 is None):
            # Se não houver espaço reservado para título, crie um textbox para o título
            title_shape1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_shape2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    
            text_frame1 = title_shape1.text_frame
            text_frame2 = title_shape2.text_frame
    
            text_frame1.text = f"{prefix}"
            text_frame2.text = f"{prefix}"
    
        else:
            title_shape1.text = f"{prefix}"
            text_frame1 = title_shape1.text_frame
            title_shape2.text = f"{prefix}"
            text_frame2 = title_shape2.text_frame
    
        for paragraph in text_frame1.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(24)  # Tamanho da fonte do título do slide
                
        for paragraph in text_frame2.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(24)  # Tamanho da fonte do título do slide
    
        for idx, grafico in enumerate(graficos):
           
            
            
            if ('100ms' in grafico or '10ms' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                # Adiciona a imagem ao slide
                # Calcula a posição da imagem
                if ('Positivo' in grafico) and ('10ms' in grafico):
                    col = 0
                    row = 0
                    left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                    top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                elif ('Positivo' in grafico) and ('100ms' in grafico):
                    col = 1
                    row = 0
                    left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                    top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                elif ('Negativo' in grafico) and ('10ms' in grafico):
                    col = 0
                    row = 1
                    left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                    top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                elif ('Negativo' in grafico) and ('100ms' in grafico):
                    col = 1
                    row = 1
                    left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                    top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
            elif ('1s' in grafico or '10s' in grafico or '40s' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                # Adiciona a imagem ao slide
                # Calcula a posição da imagem
                
    
                
                if ('Positivo' in grafico) and ('1s' in grafico):
                    
                    col = 0
                    row = 0
                    left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                    top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide2.shapes.add_picture(img_path, left, top, width=img_width_2, height=img_height_2)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                elif ('Positivo' in grafico) and ('10s' in grafico):
                    col = 1
                    row = 0
                    left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                    top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide2.shapes.add_picture(img_path, left, top, width=img_width_2, height=img_height_2)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                elif ('Positivo' in grafico) and ('40s' in grafico):
                    col = 2
                    row = 0
                    left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                    top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide2.shapes.add_picture(img_path, left, top, width=img_width_2, height=img_height_2)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                elif ('Negativo' in grafico) and ('1s' in grafico):
                    
                    col = 0
                    row = 1
                    left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                    top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide2.shapes.add_picture(img_path, left, top, width=img_width_2, height=img_height_2)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                elif ('Negativo' in grafico) and ('10s' in grafico):
                    
                    col = 1
                    row = 1
                    left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                    top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide2.shapes.add_picture(img_path, left, top, width=img_width_2, height=img_height_2)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                elif ('Negativo' in grafico) and ('40s' in grafico):
                     
                    col = 2
                    row = 1
                    left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                    top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide2.shapes.add_picture(img_path, left, top, width=img_width_2, height=img_height_2)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
    
    
    prs.save(current_directory+versionador+'Template'+versionador+'Single Pulse.pptx')
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    prs = Presentation()
    
    
    
    
    
    
    
    
    # Adiciona um slide de título
    slide_layout = prs.slide_layouts[0]  # Layout 0 é um slide de título
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Single Pulse Del IDS"
    subtitle.text = "Dados e Gráficos"
    
    # Aumenta o tamanho da fonte do título
    title_shape = title.text_frame
    for paragraph in title_shape.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(90)  # Tamanho da fonte do título
    
    # Aumenta o tamanho da fonte do subtítulo
    subtitle_shape = subtitle.text_frame
    for paragraph in subtitle_shape.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(52)  # Tamanho da fonte do subtítulo
    
    # Define o tamanho das imagens em pixels (exemplo: 300x200 pixels)
    
    
    img_width_px_1 = 320
    img_height_px_1 = 320
    
    img_width_px_2 = 320
    img_height_px_2 = 320
    
    # Converte o tamanho das imagens de pixels para polegadas
    img_width_1 = Pt(img_width_px_1 / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_1 = Pt(img_height_px_1 / 96 * 72)
    
    img_width_2 = Pt(img_width_px_2 / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_2 = Pt(img_height_px_2 / 96 * 72)
    
    # Define a posição inicial das imagens e o espaço entre elas
    margin_left_1 = Inches(1.5)
    margin_top_1 = Inches(0.8)
    horizontal_spacing_1 = Inches(0.01)
    vertical_spacing_1 = Inches(0.000)
    
    
    margin_left_2 = Inches(0.2)
    margin_top_2 = Inches(0.8)
    horizontal_spacing_2 = Inches(0.01)
    vertical_spacing_2 = Inches(0.01)
    
    
    
    # Define o caminho para a imagem
    img_caminho = 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\cnpem.png'
    
    
    img_width_px_img = 162
    img_height_px_img = 85
    
    # Converte o tamanho das imagens de pixels para polegadas
    img_width_img = Pt(img_width_px_img / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_img = Pt(img_height_px_img / 96 * 72)
    
    # Define as dimensões e a posição da imagem
    left_img = Inches(8)    # Distância da borda esquerda do slide
    top_img = Inches(0.01)     # Distância da borda superior do slide
    
    
    # Adiciona a imagem ao slide
    pic = slide.shapes.add_picture(img_caminho, left_img, top_img, img_width_img, img_height_img)
    
    # Opcional: Ajustar o tamanho da imagem sem alterar as proporções
    # pic = slide.shapes.add_picture(img_path, left, top)
    # pic.width = Inches(4)  # Define a largura
    # pic.height = Inches(3)  # Define a altura
    
    
    
    
    
    
    
    # Atualize este caminho conforme necessário
    image_folder = os.path.join(current_directory, 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\graficos_gerados\\Graficos Single Pulse\\Graficos Del Single Pulse')
    
    # Lista todas as imagens na pasta especificada
    images = [f for f in os.listdir(image_folder) if f.endswith(('png', 'jpg', 'jpeg'))]
    
    
    
    
    
    # Lista de arquivos na pasta
    files = [f for f in os.listdir(image_folder) if os.path.isfile(os.path.join(image_folder, f))]
    
    # Criação do DataFrame
    df = pd.DataFrame(files, columns=['Filename'])
    
    
    
    # Extraindo o prefixo (parte antes do '_') para usar como chave de agrupamento
    df['Group'] = df['Filename'].apply(lambda x: x.split('_')[0])
    
    
    # Agrupando por 'Group'
    grouped = df.groupby('Group')
    # Realizando o agrupamento
    
    
    # Criando um único dicionário com todos os grupos
    combined_dict = {key: group.to_dict('list') for key, group in grouped}
    
    
        
    
    
    
    
    # Criando um novo dicionário plano
    novo_dicionario = {}
    
    for chave_principal, valor in combined_dict.items():
        novo_dicionario[chave_principal] = valor['Filename']
    
    # Imprimindo o novo dicionário
    del novo_dicionario['Grafico Del Single Pulse.opju']
    print(novo_dicionario)     
    
    
    
    
    # Adiciona slides para cada grupo de gráficos
    for prefix, graficos in novo_dicionario.items():
        # Adiciona um slide vazio
        slide_layout = prs.slide_layouts[6]  # Layout 6 é um slide vazio
        slide1 = prs.slides.add_slide(slide_layout)
        pic = slide1.shapes.add_picture(img_caminho, left_img, top_img, img_width_img, img_height_img)
    
        # Adiciona título ao slide
        title_shape1 = slide1.shapes.title
        title_shape2 = slide2.shapes.title
    
        if (title_shape1 is None) or (title_shape2 is None):
            # Se não houver espaço reservado para título, crie um textbox para o título
            title_shape1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_shape2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    
            text_frame1 = title_shape1.text_frame
            text_frame2 = title_shape2.text_frame
    
            text_frame1.text = f"{prefix}"
            text_frame2.text = f"{prefix}"
    
        else:
            title_shape1.text = f"{prefix}"
            text_frame1 = title_shape1.text_frame
            title_shape2.text = f"{prefix}"
            text_frame2 = title_shape2.text_frame
    
        for paragraph in text_frame1.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(24)  # Tamanho da fonte do título do slide
                
        for paragraph in text_frame2.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(24)  # Tamanho da fonte do título do slide
    
        for idx, grafico in enumerate(graficos):
           
            
            
            
            if '0.01 s' in grafico:
                col = 0
                row = 0
                left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                print(left)
                print(top)
                img_path = os.path.join(image_folder, grafico)
                print(image_folder)
                print(grafico)
                if os.path.exists(img_path):  # Verifica se o arquivo existe
                    slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                else:
                    print(f"Arquivo não encontrado: {img_path}")
                
            elif '0.1 s' in grafico:
                col = 1
                row = 0
                left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                print(left)
                print(top)
                img_path = os.path.join(image_folder, grafico)
                print(image_folder)
                print(grafico)
                if os.path.exists(img_path):  # Verifica se o arquivo existe
                    slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                else:
                    print(f"Arquivo não encontrado: {img_path}")
            
    
            
            elif '1.0 s' in grafico:
                
                col = 0
                row = 1
                left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                print(left)
                print(top)
                img_path = os.path.join(image_folder, grafico)
                print(image_folder)
                print(grafico)
                if os.path.exists(img_path):  # Verifica se o arquivo existe
                    slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                else:
                    print(f"Arquivo não encontrado: {img_path}")
            elif '10.0 s' in grafico:
                col = 1
                row = 1
                left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                print(left)
                print(top)
                img_path = os.path.join(image_folder, grafico)
                print(image_folder)
                print(grafico)
                if os.path.exists(img_path):  # Verifica se o arquivo existe
                    slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                else:
                    print(f"Arquivo não encontrado: {img_path}")
            elif '40.0 s' in grafico:
                col = 2
                row = 1
                left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                print(left)
                print(top)
                img_path = os.path.join(image_folder, grafico)
                print(image_folder)
                print(grafico)
                if os.path.exists(img_path):  # Verifica se o arquivo existe
                    slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                else:
                    print(f"Arquivo não encontrado: {img_path}")
                
    
    
    prs.save(current_directory+versionador+'Template'+versionador+'Del IDS Single Pulse.pptx')
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    prs = Presentation()
    
    
    
    
    
    
    
    
    # Adiciona um slide de título
    slide_layout = prs.slide_layouts[0]  # Layout 0 é um slide de título
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Estabilidade"
    subtitle.text = "Dados e Gráficos"
    
    # Aumenta o tamanho da fonte do título
    title_shape = title.text_frame
    for paragraph in title_shape.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(90)  # Tamanho da fonte do título
    
    # Aumenta o tamanho da fonte do subtítulo
    subtitle_shape = subtitle.text_frame
    for paragraph in subtitle_shape.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(52)  # Tamanho da fonte do subtítulo
    
    # Define o tamanho das imagens em pixels (exemplo: 300x200 pixels)
    
    
    img_width_px_1 = 500
    img_height_px_1 = 500
    
    img_width_px_2 = 320
    img_height_px_2 = 320
    
    # Converte o tamanho das imagens de pixels para polegadas
    img_width_1 = Pt(img_width_px_1 / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_1 = Pt(img_height_px_1 / 96 * 72)
    
    img_width_2 = Pt(img_width_px_2 / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_2 = Pt(img_height_px_2 / 96 * 72)
    
    # Define a posição inicial das imagens e o espaço entre elas
    margin_left_1 = Inches(0.0)
    margin_top_1 = Inches(1.2)
    horizontal_spacing_1 = Inches(0.0)
    vertical_spacing_1 = Inches(0.000)
    
    
    margin_left_2 = Inches(0.2)
    margin_top_2 = Inches(0.8)
    horizontal_spacing_2 = Inches(0.01)
    vertical_spacing_2 = Inches(0.01)
    
    
    
    # Define o caminho para a imagem
    img_caminho = 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\cnpem.png'
    
    
    img_width_px_img = 162
    img_height_px_img = 85
    
    # Converte o tamanho das imagens de pixels para polegadas
    img_width_img = Pt(img_width_px_img / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_img = Pt(img_height_px_img / 96 * 72)
    
    # Define as dimensões e a posição da imagem
    left_img = Inches(8)    # Distância da borda esquerda do slide
    top_img = Inches(0.01)     # Distância da borda superior do slide
    
    
    # Adiciona a imagem ao slide
    pic = slide.shapes.add_picture(img_caminho, left_img, top_img, img_width_img, img_height_img)
    
    # Opcional: Ajustar o tamanho da imagem sem alterar as proporções
    # pic = slide.shapes.add_picture(img_path, left, top)
    # pic.width = Inches(4)  # Define a largura
    # pic.height = Inches(3)  # Define a altura
    
    
    
    
    
    
    
    # Atualize este caminho conforme necessário
    image_folder = os.path.join(current_directory, 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\graficos_gerados\\Graficos Estabilidade')
    
    # Lista todas as imagens na pasta especificada
    images = [f for f in os.listdir(image_folder) if f.endswith(('png', 'jpg', 'jpeg'))]
    
    
    
    
    
    # Lista de arquivos na pasta
    files = [f for f in os.listdir(image_folder) if os.path.isfile(os.path.join(image_folder, f))]
    
    # Criação do DataFrame
    df = pd.DataFrame(files, columns=['Filename'])
    df = df.drop(0).reset_index(drop=True)
    
    print(df)
    # Extraindo o prefixo (parte antes do '_') para usar como chave de agrupamento
    df['Group'] = df['Filename'].apply(lambda x: x.split('_')[1])
    print(df)
    
    
    # Agrupando por 'Group'
    grouped = df.groupby('Group')
    # Realizando o agrupamento
    
    
    # Criando um único dicionário com todos os grupos
    combined_dict = {key: group.to_dict('list') for key, group in grouped}
    
    
        
    
    
    
    
    # Criando um novo dicionário plano
    novo_dicionario = {}
    
    for chave_principal, valor in combined_dict.items():
        novo_dicionario[chave_principal] = valor['Filename']
    
    print(novo_dicionario)
    
    
    
    
    # Adiciona slides para cada grupo de gráficos
    for prefix, graficos in novo_dicionario.items():
        # Adiciona um slide vazio
        slide_layout = prs.slide_layouts[6]  # Layout 6 é um slide vazio
        slide1 = prs.slides.add_slide(slide_layout)
        pic = slide1.shapes.add_picture(img_caminho, left_img, top_img, img_width_img, img_height_img)
    
        # Adiciona título ao slide
        title_shape1 = slide1.shapes.title
        title_shape2 = slide2.shapes.title
    
        if (title_shape1 is None) or (title_shape2 is None):
            # Se não houver espaço reservado para título, crie um textbox para o título
            title_shape1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    
            text_frame1 = title_shape1.text_frame
    
            text_frame1.text = f"{prefix}"
    
        else:
            title_shape1.text = f"{prefix}"
            text_frame1 = title_shape1.text_frame
           
    
        for paragraph in text_frame1.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(36)  # Tamanho da fonte do título do slide
                
        
    
        for idx, grafico in enumerate(graficos):
           
            
            
            
            if ('KCl' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                # Adiciona a imagem ao slide
                # Calcula a posição da imagem
                if ('0.8' in grafico):
                    col = 0
                    row = 0
                    left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                    top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                elif ('0.1' in grafico):
                    col = 1
                    row = 0
                    left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                    top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
               
            elif ('MgCl2' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                # Adiciona a imagem ao slide
                # Calcula a posição da imagem
                
    
                
                if ('0.8' in grafico):
                    col = 0
                    row = 0
                    left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                    top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                elif ('0.1' in grafico):
                    col = 1
                    row = 0
                    left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                    top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                    print(left)
                    print(top)
                    img_path = os.path.join(image_folder, grafico)
                    print(image_folder)
                    print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                
            
    
            
                
    
    
    prs.save(current_directory+versionador+'Template'+versionador+'Endurance.pptx')
      

    

    prs = Presentation()
    
    
    
    
    
    
    
    
    # Adiciona um slide de título
    slide_layout = prs.slide_layouts[0]  # Layout 0 é um slide de título
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Tempo de Retenção"
    subtitle.text = "Dados e Gráficos"

    # Aumenta o tamanho da fonte do título
    title_shape = title.text_frame
    for paragraph in title_shape.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(90)  # Tamanho da fonte do título

    # Aumenta o tamanho da fonte do subtítulo
    subtitle_shape = subtitle.text_frame
    for paragraph in subtitle_shape.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(52)  # Tamanho da fonte do subtítulo

    # Define o tamanho das imagens em pixels (exemplo: 300x200 pixels)
    
    
    img_width_px_1 = 320
    img_height_px_1 = 320
    
    img_width_px_2 = 320
    img_height_px_2 = 320
    
    img_width_px_3 = 600
    img_height_px_3 = 600
    
    # Converte o tamanho das imagens de pixels para polegadas
    img_width_1 = Pt(img_width_px_1 / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_1 = Pt(img_height_px_1 / 96 * 72)
    
    img_width_2 = Pt(img_width_px_2 / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_2 = Pt(img_height_px_2 / 96 * 72)
    
    img_width_3 = Pt(img_width_px_3 / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_3 = Pt(img_height_px_3 / 96 * 72)
    
    # Define a posição inicial das imagens e o espaço entre elas
    margin_left_1 = Inches(0.0)
    margin_top_1 = Inches(2.2)
    horizontal_spacing_1 = Inches(0.0)
    vertical_spacing_1 = Inches(0.000)
    
    
    margin_left_2 = Inches(0.0)
    margin_top_2 = Inches(0.8)
    horizontal_spacing_2 = Inches(0.01)
    vertical_spacing_2 = Inches(0.01)
    
    
    margin_left_3 = Inches(1.8)
    margin_top_3 = Inches(0.8)
    horizontal_spacing_3 = Inches(0.01)
    vertical_spacing_3 = Inches(0.01)
    
    
    # Define o caminho para a imagem
    img_caminho = 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\cnpem.png'
    
    
    img_width_px_img = 162
    img_height_px_img = 85
    
    # Converte o tamanho das imagens de pixels para polegadas
    img_width_img = Pt(img_width_px_img / 96 * 72)  # 1 polegada = 96 pixels (resolução padrão) e 1 polegada = 72 pontos
    img_height_img = Pt(img_height_px_img / 96 * 72)
    
    # Define as dimensões e a posição da imagem
    left_img = Inches(8)    # Distância da borda esquerda do slide
    top_img = Inches(0.01)     # Distância da borda superior do slide
    
    
    # Adiciona a imagem ao slide
    pic = slide.shapes.add_picture(img_caminho, left_img, top_img, img_width_img, img_height_img)
    
    # Opcional: Ajustar o tamanho da imagem sem alterar as proporções
    # pic = slide.shapes.add_picture(img_path, left, top)
    # pic.width = Inches(4)  # Define a largura
    # pic.height = Inches(3)  # Define a altura
    
    
    
    
    
    

        
    # Lista de caminhos de pastas
    image_folders = [
        os.path.join(current_directory, 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\graficos_gerados\\Graficos Transfer Tempo de Retencao'),
        # Adicione mais caminhos conforme necessário
        os.path.join(current_directory, 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\graficos_gerados\\Graficos Transcondutancia'),
        os.path.join(current_directory, 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\graficos_gerados\\Graficos Pulso Longo'),
        os.path.join(current_directory, 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\graficos_gerados\\Grafico IGS'),
        os.path.join(current_directory, 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\graficos_gerados\\Grafico Trans Short Pulse'),
        os.path.join(current_directory, 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\graficos_gerados\\Grafico Trans Short Pulse\\Grafico Transcondutancia Short Pulse'),
        os.path.join(current_directory, 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\graficos_gerados\\Grafico Trans Short Pulse\\Grafico IGS Short Pulse'),
        os.path.join(current_directory, 'C:\\Users\\eduardo.neto\\Desktop\\programa_v5\\graficos_gerados\\Graficos Pulso Curto')
    ]
    
    # Inicializa listas para armazenar todas as imagens e arquivos
    all_image_paths = []
    all_file_paths = []
    
    # Itera sobre cada pasta na lista
    for folder in image_folders:
        # Verifica se o caminho da pasta existe
        if os.path.exists(folder):
            # Lista todas as imagens na pasta atual e armazena o caminho completo
            image_paths = [os.path.join(folder, f) for f in os.listdir(folder) if f.endswith(('png', 'jpg', 'jpeg'))]
            all_image_paths.extend(image_paths)
            
            # Lista de arquivos na pasta atual e armazena o caminho completo
            file_paths = [os.path.join(folder, f) for f in os.listdir(folder) if os.path.isfile(os.path.join(folder, f))]
            all_file_paths.extend(file_paths)
        else:
            print(f"A pasta {folder} não existe.")
        
    # Exibe todas as imagens e arquivos encontrados
    print(all_file_paths)
    print("AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA")
    print(all_image_paths)





    # Criação do DataFrame
    df = pd.DataFrame(all_file_paths, columns=['Filename'])
    
    # Filtra o DataFrame para manter apenas os caminhos que levam a imagens válidas
    valid_image_extensions = ('.png', '.jpg', '.jpeg')
    df = df[df['Filename'].str.endswith(valid_image_extensions)].reset_index(drop=True)
    
    # Exibe o DataFrame filtrado
    print(df)
    
    # Extraindo o prefixo (parte antes do '_') para usar como chave de agrupamento
    df['Group'] = df['Filename'].apply(lambda x: x.split('&')[1])
    print("Aqui deu certoooooooooooooooooooooooooo\n\n\n\n\n\n\n\n")
    print(df)
    

    
    # Agrupando por 'Group'
    grouped = df.groupby('Group')
    # Realizando o agrupamento
    
    
    # Criando um único dicionário com todos os grupos
    combined_dict = {key: group.to_dict('list') for key, group in grouped}
    
    
        
    
    
    

    # Criando um novo dicionário plano
    novo_dicionario = {}

    for chave_principal, valor in combined_dict.items():
        novo_dicionario[chave_principal] = valor['Filename']
    # Imprimindo o novo dicionário
    print(novo_dicionario) 




    # Adiciona slides para cada grupo de gráficos
    for prefix, graficos in novo_dicionario.items():
        # Adiciona um slide vazio
        slide_layout = prs.slide_layouts[6]  # Layout 6 é um slide vazio
        slide1 = prs.slides.add_slide(slide_layout)
        slide2 = prs.slides.add_slide(slide_layout)
        slide3 = prs.slides.add_slide(slide_layout)
        pic = slide3.shapes.add_picture(img_caminho, left_img, top_img, img_width_img, img_height_img)
        pic = slide2.shapes.add_picture(img_caminho, left_img, top_img, img_width_img, img_height_img)
        pic = slide1.shapes.add_picture(img_caminho, left_img, top_img, img_width_img, img_height_img)
        

        # Adiciona título ao slide
        title_shape1 = slide1.shapes.title
        title_shape2 = slide2.shapes.title
        title_shape3 = slide3.shapes.title

        if (title_shape1 is None) or (title_shape2 is None):
            # Se não houver espaço reservado para título, crie um textbox para o título
            title_shape1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_shape2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_shape3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))


            text_frame1 = title_shape1.text_frame
            text_frame2 = title_shape2.text_frame
            text_frame3 = title_shape3.text_frame


            text_frame1.text = f"{prefix[:-4]}"
            text_frame2.text = f"{prefix[:-4]}"
            text_frame3.text = f"{prefix[:-4]}"


        else:
            title_shape1.text = f"{prefix[:-4]}"
            text_frame1 = title_shape1.text_frame
            title_shape2.text = f"{prefix[:-4]}"
            text_frame2 = title_shape2.text_frame
            title_shape3.text = f"{prefix[:-4]}"
            text_frame3 = title_shape3.text_frame

        for paragraph in text_frame1.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(24)  # Tamanho da fonte do título do slide
                
        for paragraph in text_frame2.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(24)  # Tamanho da fonte do título do slide
                
        for paragraph in text_frame3.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(24)  # Tamanho da fonte do título do slide

        for idx, grafico in enumerate(graficos):
           
            
            
            if ('Retention Transfer' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                # Adiciona a imagem ao slide
                # Calcula a posição da imagem
              
                    col = 0
                    row = 0
                    left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                    top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                    #print(left)
                   # print(top)
                    img_path = os.path.join(grafico)
                    print(folder)
                    #print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
            elif ('Transconduct&' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                    col = 1
                    row = 0
                    left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                    top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                   # print(left)
                    #print(top)
                    img_path = os.path.join(grafico)
                    print(folder)
                    #print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
            elif ('IGS&' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                    col = 2
                    row = 0
                    left = margin_left_1 + col * (img_width_1 + horizontal_spacing_1)
                    top = margin_top_1 + row * (img_height_1 + vertical_spacing_1)
                   # print(left)
                    #print(top)
                    img_path = os.path.join(grafico)
                    print(folder)
                    #print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide1.shapes.add_picture(img_path, left, top, width=img_width_1, height=img_height_1)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                        
                        
                        
                        
                        
            
            
            
            
            elif ('IGS Pulso Curto&' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                    col = 0
                    row = 0
                    left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                    top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                   # print(left)
                    #print(top)
                    img_path = os.path.join(grafico)
                    print(folder)
                    #print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide2.shapes.add_picture(img_path, left, top, width=img_width_2, height=img_height_2)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")          
            elif ('Pulso Curto&' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                    col = 1
                    row = 1
                    left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                    top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                   # print(left)
                    #print(top)
                    img_path = os.path.join(grafico)
                    print(folder)
                    #print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide2.shapes.add_picture(img_path, left, top, width=img_width_2, height=img_height_2)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
            
            elif ('Transfer Short Pulse' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                    col = 2
                    row = 0
                    left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                    top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                   # print(left)
                    #print(top)
                    img_path = os.path.join(grafico)
                    print(folder)
                    #print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide2.shapes.add_picture(img_path, left, top, width=img_width_2, height=img_height_2)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                        
            elif ('IGS Short Pulse&' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                    col = 2
                    row = 1
                    left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                    top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                   # print(left)
                    #print(top)
                    img_path = os.path.join(grafico)
                    print(folder)
                    #print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide2.shapes.add_picture(img_path, left, top, width=img_width_2, height=img_height_2)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
            elif ('Transconduct Short Pulse&' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                    col = 1
                    row = 0
                    left = margin_left_2 + col * (img_width_2 + horizontal_spacing_2)
                    top = margin_top_2 + row * (img_height_2 + vertical_spacing_2)
                   # print(left)
                    #print(top)
                    img_path = os.path.join(grafico)
                    print(folder)
                    #print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide2.shapes.add_picture(img_path, left, top, width=img_width_2, height=img_height_2)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                
                    
                

                
            elif ('Pulso Longo' in grafico) and grafico.endswith(('png', 'jpg', 'jpeg')):
                    
                    col = 0
                    row = 0
                    left = margin_left_3 + col * (img_width_3 + horizontal_spacing_3)
                    top = margin_top_3 + row * (img_height_3 + vertical_spacing_3)
                   # print(left)
                    #print(top)
                    img_path = os.path.join(grafico)
                    print(folder)
                    #print(grafico)
                    if os.path.exists(img_path):  # Verifica se o arquivo existe
                        slide3.shapes.add_picture(img_path, left, top, width=img_width_3, height=img_height_3)
                    else:
                        print(f"Arquivo não encontrado: {img_path}")
                
               
                
                
              


    prs.save(current_directory+versionador+'Template'+versionador+'tempo de Retenção.pptx')








    



